import json
import pandas as pd
from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from models import SlideFieldSelection
from database import async_session, Project, Slide
from sqlalchemy import select


class PowerPointService:
    """Service for generating and updating PowerPoint presentations"""
    
    def __init__(self, project_id: str):
        self.project_id = project_id
    
    async def generate_complete_presentation(self) -> str:
        """Generate complete PowerPoint presentation with all slides"""
        try:
            async with async_session() as session:
                # Get project data
                result = await session.execute(
                    select(Project).where(Project.id == self.project_id)
                )
                project = result.scalar_one()
                
                # Get all slides for this project
                slides_result = await session.execute(
                    select(Slide).where(Slide.project_id == self.project_id)
                    .order_by(Slide.slide_number)
                )
                slides = slides_result.scalars().all()
                
                # Load data
                data_df = pd.read_csv(project.data_source_path)
                
                # Create presentation
                if project.template_path and Path(project.template_path).exists():
                    prs = Presentation(project.template_path)
                else:
                    prs = Presentation()
                
                # Ensure we have enough slides in the presentation
                while len(prs.slides) < len(slides):
                    # Add a blank slide with title and content layout
                    slide_layout = prs.slide_layouts[1]  # Title and Content layout
                    prs.slides.add_slide(slide_layout)
                
                # Update each slide with data
                for i, slide_data in enumerate(slides):
                    if i < len(prs.slides):
                        ppt_slide = prs.slides[i]
                        await self._update_slide_content(
                            ppt_slide, slide_data, data_df, i + 1
                        )
                
                # Save presentation
                output_dir = Path(f"downloads/{self.project_id}")
                output_dir.mkdir(parents=True, exist_ok=True)
                output_path = output_dir / "analysis_report.pptm"
                
                prs.save(str(output_path))
                return str(output_path)
                
        except Exception as e:
            raise Exception(f"Failed to generate PowerPoint: {str(e)}")
    
    async def _update_slide_content(
        self, ppt_slide: Any, slide_data: Slide, data_df: pd.DataFrame, slide_number: int
    ):
        """Update a single slide with data table"""
        try:
            # Set slide title
            if hasattr(ppt_slide, 'shapes') and len(ppt_slide.shapes) > 0:
                title_shape = ppt_slide.shapes.title
                if title_shape and hasattr(title_shape, 'text'):
                    title_shape.text = slide_data.slide_title or f"Slide {slide_number}"
            
            # Get fields to display (use final_fields if available, otherwise agent_selected_fields)
            fields_data = slide_data.final_fields or slide_data.agent_selected_fields
            if not fields_data:
                return
            
            # Convert dict to SlideFieldSelection objects
            fields = [SlideFieldSelection(**field) for field in fields_data]
            
            # Find content placeholder or create table area
            content_placeholder = None
            for shape in ppt_slide.shapes:
                if hasattr(shape, 'text') and shape != ppt_slide.shapes.title:
                    content_placeholder = shape
                    break
            
            # Remove existing content placeholder if found
            if content_placeholder:
                shape_element = content_placeholder.element
                shape_element.getparent().remove(shape_element)
            
            # Create table
            await self._create_data_table(ppt_slide, fields, data_df)
            
        except Exception as e:
            print(f"Error updating slide {slide_number}: {e}")
    
    async def _create_data_table(
        self, slide: Any, fields: List[SlideFieldSelection], data_df: pd.DataFrame
    ):
        """Create a data table on the slide based on field configuration"""
        try:
            # Get all unique metric fields across all rows (these become columns)
            all_metric_fields = []
            for field in fields:
                if field.metric_fields:
                    for metric in field.metric_fields:
                        if metric not in all_metric_fields:
                            all_metric_fields.append(metric)
            
            if not all_metric_fields:
                return
            
            # Calculate table dimensions
            rows = len(fields) + 1  # +1 for header
            cols = len(all_metric_fields) + 1  # +1 for row label column
            
            # Add table to slide
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Style table
            table.table_style = None  # Remove default style
            
            # Set header row
            self._set_table_headers_new(table, all_metric_fields)
            
            # Fill data rows
            for i, field in enumerate(fields):
                row_idx = i + 1
                self._fill_table_row_new(table, row_idx, field, all_metric_fields, data_df)
            
            # Apply formatting
            self._format_table_new(table, fields)
            
        except Exception as e:
            print(f"Error creating table: {e}")
    
    def _set_table_headers_new(self, table: Any, all_metric_fields: List[str]):
        """Set table header row with all metric fields as columns"""
        try:
            # First column is empty (for row labels, no header text)
            first_cell = table.cell(0, 0)
            first_cell.text = ""
            self._style_header_cell(first_cell)
            
            # Rest of columns are metric fields
            for i, metric_field in enumerate(all_metric_fields):
                col_idx = i + 1
                cell = table.cell(0, col_idx)
                cell.text = self._format_header_name(metric_field)
                self._style_header_cell(cell)
                
        except Exception as e:
            print(f"Error setting table headers: {e}")
    
    def _fill_table_row_new(
        self, table: Any, row_idx: int, field: SlideFieldSelection, 
        all_metric_fields: List[str], data_df: pd.DataFrame
    ):
        """Fill a table row with data based on field configuration"""
        try:
            if row_idx >= table.rows:
                return
            
            # Set row label (first column)
            label_cell = table.cell(row_idx, 0)
            label_cell.text = field.row_label
            
            if field.is_group_header:
                # Group header styling
                self._style_group_header_cell(label_cell)
                
                if getattr(field, 'spans_all_columns', False):
                    # Group header spans all columns - clear other cells
                    for col_idx in range(1, table.columns):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = ""
                        self._style_group_header_cell(cell)
                else:
                    # Fill metric columns for group header
                    self._fill_metric_columns(table, row_idx, field, all_metric_fields, data_df, is_group_header=True)
            else:
                # Regular data row
                self._style_data_cell(label_cell)
                self._fill_metric_columns(table, row_idx, field, all_metric_fields, data_df, is_group_header=False)
                        
        except Exception as e:
            print(f"Error filling table row: {e}")
    
    def _fill_metric_columns(
        self, table: Any, row_idx: int, field: SlideFieldSelection,
        all_metric_fields: List[str], data_df: pd.DataFrame, is_group_header: bool = False
    ):
        """Fill metric columns for a row"""
        try:
            for col_idx, metric_field in enumerate(all_metric_fields):
                cell = table.cell(row_idx, col_idx + 1)
                
                if metric_field in (field.metric_fields or []):
                    # This row uses this metric field - calculate value
                    if metric_field in data_df.columns:
                        value = self._calculate_metric_value_with_filters(
                            data_df, metric_field, field.aggregation,
                            getattr(field, 'filters', [])
                        )
                        cell.text = self._format_value(value, metric_field)
                    else:
                        # Field not found in data
                        cell.text = "N/A"
                        print(f"Warning: Field '{metric_field}' not found in data. Available fields: {list(data_df.columns)}")
                    
                    # Apply styling
                    if is_group_header:
                        self._style_group_header_cell(cell)
                    else:
                        self._style_data_cell(cell)
                else:
                    # This row doesn't use this metric - empty cell
                    cell.text = ""
                    if is_group_header:
                        self._style_group_header_cell(cell)
                    else:
                        self._style_data_cell(cell)
                        
        except Exception as e:
            print(f"Error filling metric columns: {e}")
    
    def _calculate_metric_value_with_filters(
        self, data_df: pd.DataFrame, metric_field: str, aggregation: str, filters: List[Dict]
    ) -> float:
        """Calculate aggregated value for a metric field with optional filters"""
        try:
            # Start with the full dataframe
            filtered_df = data_df.copy()
            
            # Apply filters if any
            if filters:
                for filter_condition in filters:
                    if isinstance(filter_condition, dict):
                        field_name = filter_condition.get('field')
                        operator = filter_condition.get('operator', '==')
                        value = filter_condition.get('value')
                        
                        if field_name and field_name in filtered_df.columns:
                            if operator == '==':
                                filtered_df = filtered_df[filtered_df[field_name] == value]
                            elif operator == '!=':
                                filtered_df = filtered_df[filtered_df[field_name] != value]
                            elif operator == '>':
                                filtered_df = filtered_df[filtered_df[field_name] > value]
                            elif operator == '<':
                                filtered_df = filtered_df[filtered_df[field_name] < value]
                            elif operator == '>=':
                                filtered_df = filtered_df[filtered_df[field_name] >= value]
                            elif operator == '<=':
                                filtered_df = filtered_df[filtered_df[field_name] <= value]
            
            # Check if metric field exists
            if metric_field not in filtered_df.columns:
                return 0
            
            # Apply aggregation
            if aggregation == "sum":
                return filtered_df[metric_field].sum()
            elif aggregation == "average" or aggregation == "mean":
                return filtered_df[metric_field].mean()
            elif aggregation == "count":
                return filtered_df[metric_field].count()
            elif aggregation == "max":
                return filtered_df[metric_field].max()
            elif aggregation == "min":
                return filtered_df[metric_field].min()
            else:
                return filtered_df[metric_field].sum()  # Default to sum
                
        except Exception as e:
            print(f"Error calculating metric value: {e}")
            return 0
    
    def _format_table_new(self, table: Any, fields: List[SlideFieldSelection]):
        """Apply overall table formatting with enhanced styling"""
        try:
            # Set table borders and general styling
            for row_idx, row in enumerate(table.rows):
                row.height = Inches(0.5)
                
                # Adjust height for group headers
                if row_idx > 0:  # Skip header row
                    field_idx = row_idx - 1
                    if field_idx < len(fields) and fields[field_idx].is_group_header:
                        row.height = Inches(0.6)  # Slightly taller for group headers
            
            for col in table.columns:
                col.width = Inches(1.5)
            
            # Make first column wider for labels
            if table.columns:
                table.columns[0].width = Inches(3.0)
                
        except Exception as e:
            print(f"Error formatting table: {e}")
    

    
    def _format_header_name(self, field_name: str) -> str:
        """Format field name for display"""
        return field_name.replace("_", " ").title()
    
    def _style_header_cell(self, cell: Any):
        """Apply styling to header cells"""
        try:
            # Set font
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(12)
                    run.font.color.rgb = RGBColor(255, 255, 255)
                
                paragraph.alignment = PP_ALIGN.CENTER
            
            # Set background color (dark blue)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(31, 78, 121)
            
        except Exception as e:
            print(f"Error styling header cell: {e}")
    
    def _style_group_header_cell(self, cell: Any):
        """Apply styling to group header cells"""
        try:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                
                paragraph.alignment = PP_ALIGN.LEFT
            
            # Set background color (light gray)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(242, 242, 242)
            
        except Exception as e:
            print(f"Error styling group header cell: {e}")
    
    def _style_data_cell(self, cell: Any):
        """Apply styling to data cells"""
        try:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                
                paragraph.alignment = PP_ALIGN.RIGHT
            
        except Exception as e:
            print(f"Error styling data cell: {e}")
    
    def _format_value(self, value: float, field_name: str) -> str:
        """Format value based on field type and context"""
        try:
            if value == 0:
                return "-"
                
            # Insurance/actuarial specific formatting
            if any(keyword in field_name.lower() for keyword in ['actualincurred', 'nominalreserves', 'discountedreserves', 'ocl', 'changeinocl', 'reserves', 'incurred', 'claim']):
                return f"${value:,.0f}"
            elif any(keyword in field_name.lower() for keyword in ['rate', 'ratio', 'percent']):
                return f"{value:.2%}" if value <= 1 else f"{value:.2f}%"
            elif any(keyword in field_name.lower() for keyword in ['count', 'number', 'quantity', 'year']):
                return f"{value:,.0f}"
            elif isinstance(value, float):
                return f"{value:,.2f}"
            else:
                return str(value)
        except Exception:
            return "N/A"
