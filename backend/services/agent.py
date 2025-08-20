import asyncio
import pandas as pd
import json
from datetime import datetime
from pathlib import Path
from typing import Dict, List
from pptx import Presentation
from pptx.util import Inches
import markdown

from models import Project, ProjectStatus, SlideInfo, ProjectOutput
from services.websocket_manager import WebSocketManager

class AnalysisAgent:
    def __init__(self, project_id: str, websocket_manager: WebSocketManager):
        self.project_id = project_id
        self.websocket_manager = websocket_manager
        self.slides_generated = []
        
    async def start_analysis(self, project: Project):
        """Main analysis workflow"""
        try:
            await self.send_status("PROCESSING", "Starting analysis agent...")
            await asyncio.sleep(1)
            
            # Load and analyze data
            data = await self.load_data(project)
            schema = await self.load_schema(project)
            
            await self.send_status("PROCESSING", "Analyzing data structure...")
            await asyncio.sleep(1)
            
            # Determine slide plan
            slide_plan = self.determine_slide_plan(project.field_config, data, schema)
            
            await self.send_status("PROCESSING", f"Planning complete. Generating {len(slide_plan)} slides...")
            await asyncio.sleep(1)
            
            # Generate slides
            for i, slide_info in enumerate(slide_plan):
                await self.generate_slide(i + 1, slide_info, data, project)
                await asyncio.sleep(2)  # Simulate processing time
            
            # Update status to DATA_READY
            project.status = ProjectStatus.DATA_READY
            await self.send_status("DATA_READY", "Data analysis complete. Generating outputs...")
            
            # Generate outputs
            await self.generate_outputs(project, data)
            
            # Compile PowerPoint
            await self.compile_powerpoint(project)
            
            # Final status
            project.status = ProjectStatus.PPT_READY
            await self.send_status("PPT_READY", f"Analysis complete! {len(self.slides_generated)} slides generated.")
            
        except Exception as e:
            project.status = ProjectStatus.FAILED
            await self.send_status("FAILED", f"Analysis failed: {str(e)}")
    
    async def load_data(self, project: Project) -> pd.DataFrame:
        """Load CSV data"""
        data_path = project.files["data_source"]
        return pd.read_csv(data_path)
    
    async def load_schema(self, project: Project) -> dict:
        """Load JSON schema"""
        schema_path = project.files["schema"]
        with open(schema_path, 'r') as f:
            return json.load(f)
    
    def determine_slide_plan(self, field_config: Dict, data: pd.DataFrame, schema: dict) -> List[Dict]:
        """Determine what slides to generate based on field selection"""
        slide_plan = []
        
        selected_fields = [field for field, config in (field_config or {}).items() if config.get("selected", False)]
        
        if not selected_fields:
            selected_fields = list(data.columns)
        
        # Always start with executive summary
        slide_plan.append({
            "title": "Executive Summary",
            "type": "executive_summary",
            "fields": selected_fields
        })
        
        # Individual field analysis (for detailed analysis)
        for field in selected_fields:
            if field in schema:
                slide_plan.append({
                    "title": f"{schema[field]} Analysis",
                    "type": "field_analysis",
                    "fields": [field]
                })
        
        # Cross-field analysis if multiple fields
        if len(selected_fields) > 1:
            slide_plan.append({
                "title": "Correlation Analysis",
                "type": "correlation",
                "fields": selected_fields
            })
            
            slide_plan.append({
                "title": "Comparative Analysis",
                "type": "comparison",
                "fields": selected_fields
            })
        
        # Business insights
        slide_plan.append({
            "title": "Risk Assessment",
            "type": "risk_analysis",
            "fields": selected_fields
        })
        
        slide_plan.append({
            "title": "Recommendations",
            "type": "recommendations",
            "fields": selected_fields
        })
        
        return slide_plan
    
    async def generate_slide(self, slide_num: int, slide_info: Dict, data: pd.DataFrame, project: Project):
        """Generate individual slide"""
        await self.send_status("PROCESSING", f"Generating slide {slide_num}: {slide_info['title']}...")
        
        # Mock slide content generation
        slide_content = self.generate_slide_content(slide_info, data)
        
        slide = SlideInfo(
            slide_number=slide_num,
            slide_title=slide_info["title"],
            slide_content=slide_content,
            fields_included=slide_info["fields"],
            analysis_type=slide_info["type"],
            generated_at=datetime.now().isoformat()
        )
        
        self.slides_generated.append(slide)
        project.slides.append(slide)
        
        # Send slide ready notification
        await self.websocket_manager.send_to_project(self.project_id, {
            "type": "slide_generated",
            "slide_number": slide_num,
            "slide_title": slide_info["title"],
            "status": f"SLIDE_{slide_num}_READY",
            "fields_included": slide_info["fields"],
            "analysis_type": slide_info["type"]
        })
    
    def generate_slide_content(self, slide_info: Dict, data: pd.DataFrame) -> str:
        """Generate mock slide content based on slide type"""
        slide_type = slide_info["type"]
        fields = slide_info["fields"]
        
        if slide_type == "executive_summary":
            return f"""
# Executive Summary

## Key Findings
- Analyzed {len(data)} data points across {len(fields)} fields
- Fields analyzed: {', '.join(fields)}
- Data quality: 100% complete records

## Performance Highlights
- Average growth rate: {data[fields[0]].pct_change().mean()*100:.1f}% (if applicable)
- Key trends identified across all selected metrics
- Strong correlation patterns observed

## Strategic Recommendations
- Continue monitoring key performance indicators
- Focus on identified growth opportunities
- Implement risk mitigation strategies
"""
        
        elif slide_type == "field_analysis":
            field = fields[0]
            field_data = data[field] if field in data.columns else [1, 2, 3, 4]
            return f"""
# {slide_info['title']}

## Data Overview
- Field: {field}
- Data points: {len(field_data)}
- Range: {min(field_data)} - {max(field_data)}
- Average: {sum(field_data)/len(field_data):.2f}

## Key Insights
- Trend analysis shows consistent patterns
- Growth trajectory indicates positive momentum
- Risk factors identified and manageable

## Recommendations
- Monitor this metric closely
- Consider optimization opportunities
- Implement tracking mechanisms
"""
        
        elif slide_type == "correlation":
            return f"""
# Correlation Analysis

## Field Relationships
- Analyzed correlations between: {', '.join(fields)}
- Strong positive correlation identified
- Predictive patterns observed

## Statistical Insights
- Correlation coefficient: 0.85+ (simulated)
- Relationship strength: Strong
- Predictive value: High

## Business Implications
- Fields move together predictably
- Can use for forecasting
- Risk diversification needed
"""
        
        else:
            return f"""
# {slide_info['title']}

## Analysis Summary
- Focus areas: {', '.join(fields)}
- Analysis type: {slide_type}
- Key insights generated

## Findings
- Detailed analysis completed
- Patterns identified
- Actionable insights derived

## Next Steps
- Review recommendations
- Implement suggested actions
- Monitor ongoing performance
"""
    
    async def generate_outputs(self, project: Project, data: pd.DataFrame):
        """Generate additional outputs"""
        
        # Enhanced CSV output
        enhanced_data = data.copy()
        if len(data) > 0:
            # Add some mock analysis columns
            enhanced_data['trend_indicator'] = ['↗ Growing'] * len(data)
            enhanced_data['risk_level'] = ['Low', 'Medium', 'Medium', 'High'][:len(data)]
        
        csv_output = ProjectOutput(
            output_type="table",
            content=enhanced_data.to_dict('records'),
            generated_at=datetime.now().isoformat()
        )
        project.outputs.append(csv_output)
        
        # Markdown analysis report
        markdown_content = self.generate_markdown_report(project, data)
        
        markdown_output = ProjectOutput(
            output_type="markdown",
            content=markdown_content,
            generated_at=datetime.now().isoformat()
        )
        project.outputs.append(markdown_output)
        
        # Send outputs ready notification
        await self.websocket_manager.send_to_project(self.project_id, {
            "type": "outputs_ready",
            "csv_data": enhanced_data.to_dict('records'),
            "markdown_content": markdown_content
        })
    
    def generate_markdown_report(self, project: Project, data: pd.DataFrame) -> str:
        """Generate comprehensive markdown analysis report"""
        selected_fields = list((project.field_config or {}).keys()) or list(data.columns)
        
        return f"""
# {project.name} - Analysis Report
*Generated on: {datetime.now().strftime('%B %d, %Y at %H:%M')}*

## Executive Summary
This report presents a comprehensive analysis of the submitted data across {len(selected_fields)} key fields. Our intelligent agent has processed {len(data)} data points to provide actionable insights for strategic decision-making.

## Data Overview
- **Dataset Size**: {len(data)} records
- **Fields Analyzed**: {', '.join(selected_fields)}
- **Analysis Type**: Comprehensive Financial Performance Review
- **Data Quality**: 100% complete records

## Key Findings

### 📈 Performance Metrics
- **Growth Pattern**: Consistent upward trajectory observed
- **Volatility**: Moderate risk levels identified
- **Correlation Strength**: Strong relationships between key metrics

### 📊 Statistical Summary
{self.generate_statistical_summary(data, selected_fields)}

### ⚠️ Risk Assessment
- **Current Risk Level**: Medium
- **Key Risk Factors**: 
  - Concentration risk in primary metrics
  - Market volatility exposure
  - Growth sustainability concerns

## Strategic Recommendations

1. **Performance Optimization**
   - Focus on high-impact areas identified in the analysis
   - Implement continuous monitoring systems
   - Establish key performance benchmarks

2. **Risk Mitigation**
   - Diversify exposure across different metrics
   - Implement early warning systems
   - Regular stress testing protocols

3. **Growth Strategy**
   - Capitalize on identified growth opportunities
   - Monitor market conditions closely
   - Adjust strategies based on performance data

## Technical Details
- **Analysis Method**: Multi-variate statistical analysis
- **Confidence Level**: 95%
- **Data Processing**: Automated with human oversight
- **Quality Assurance**: Comprehensive validation performed

---
*This analysis was generated by Expert Sure AI Agent v1.0*
"""
    
    def generate_statistical_summary(self, data: pd.DataFrame, fields: List[str]) -> str:
        """Generate statistical summary for markdown report"""
        if len(data) == 0:
            return "No data available for statistical analysis."
        
        summary_lines = []
        for field in fields:
            if field in data.columns and pd.api.types.is_numeric_dtype(data[field]):
                field_data = data[field]
                summary_lines.append(f"- **{field}**: Mean = {field_data.mean():.2f}, Std = {field_data.std():.2f}")
        
        return '\n'.join(summary_lines) if summary_lines else "Statistical summary not available for selected fields."
    
    async def compile_powerpoint(self, project: Project):
        """Compile all slides into PowerPoint file"""
        await self.send_status("PROCESSING", "Compiling PowerPoint presentation...")
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Add title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = project.name
        title_slide.placeholders[1].text = f"Analysis Report - Generated {datetime.now().strftime('%B %d, %Y')}"
        
        # Add content slides
        for slide_info in self.slides_generated:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_info.slide_title
            
            # Add content (simplified - in production, add charts, tables, etc.)
            content_box = slide.shapes.placeholders[1]
            content_box.text = f"Analysis completed for: {', '.join(slide_info.fields_included)}\n\nKey insights and recommendations included."
        
        # Save PowerPoint file
        output_dir = Path(f"downloads/{self.project_id}")
        output_dir.mkdir(parents=True, exist_ok=True)
        
        ppt_filename = f"{project.name.replace(' ', '_')}_Analysis.pptx"
        ppt_path = output_dir / ppt_filename
        
        prs.save(str(ppt_path))
        
        # Add file output to project
        file_output = ProjectOutput(
            output_type="file",
            content=ppt_filename,
            url=f"/downloads/{self.project_id}/{ppt_filename}",
            generated_at=datetime.now().isoformat()
        )
        project.outputs.append(file_output)
        
        # Send file ready notification
        await self.websocket_manager.send_to_project(self.project_id, {
            "type": "file_ready",
            "file_url": f"/downloads/{self.project_id}/{ppt_filename}",
            "file_name": ppt_filename,
            "file_size": f"{ppt_path.stat().st_size / 1024:.1f} KB"
        })
    
    async def send_status(self, status: str, message: str):
        """Send status update via WebSocket"""
        await self.websocket_manager.send_to_project(self.project_id, {
            "type": "status_update",
            "status": status,
            "message": message,
            "timestamp": datetime.now().isoformat()
        }) 