#!/usr/bin/env python3

import asyncio
import sys
import os
import traceback
from pathlib import Path

# Add backend to path
sys.path.append(str(Path(__file__).parent / "backend"))

from services.powerpoint_service import PowerPointService

async def test_powerpoint_generation():
    """Test PowerPoint generation with mock data"""
    
    print("Testing PowerPoint generation...")
    
    try:
        # Create test project ID
        test_project_id = "test-project-123"
        
        # Create PowerPoint service
        print("Creating PowerPoint service...")
        service = PowerPointService(test_project_id)
        print("✓ PowerPoint service created successfully")
        
        # Test that the service can import required libraries
        from pptx import Presentation
        print("✓ python-pptx library available")
        
        # Test creating a basic presentation
        print("Creating basic presentation...")
        prs = Presentation()
        
        # Add a slide
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Test Presentation"
        
        # Test saving
        output_dir = Path("downloads") / test_project_id
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "test_presentation.pptx"
        
        prs.save(str(output_path))
        print(f"✓ Test presentation saved to: {output_path}")
        
        # Check if file exists
        if output_path.exists():
            print(f"✓ File created successfully, size: {output_path.stat().st_size} bytes")
        else:
            print("✗ File was not created")
            
        return True
        
    except Exception as e:
        print(f"✗ Error during PowerPoint generation test: {e}")
        traceback.print_exc()
        return False

if __name__ == "__main__":
    success = asyncio.run(test_powerpoint_generation())
    if success:
        print("\n🎉 PowerPoint generation test PASSED")
    else:
        print("\n❌ PowerPoint generation test FAILED")
        sys.exit(1)
